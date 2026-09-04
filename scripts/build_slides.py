#!/usr/bin/env python3
"""Build the Session 0 lecture deck from one content definition.

    uv run --with python-pptx --with matplotlib --with numpy scripts/build_slides.py

Outputs (all committed):
    public/slides/img/*.png                                figures, computed from the car table
    public/slides/session-0-least-squares-foundations.pptx editable PowerPoint (python-pptx)
    public/slides/session-0-least-squares-foundations.md   the same deck as Marp markdown

Everything numeric on the slides (means, correlation, eigenvalues, the least squares fit)
is computed here from the five cars, so the numbers cannot drift from the figures.
"""
from __future__ import annotations

import pathlib

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.patches import Ellipse, FancyArrowPatch, Polygon
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = pathlib.Path(__file__).resolve().parent.parent
OUT = ROOT / "public" / "slides"
IMG = OUT / "img"
IMG.mkdir(parents=True, exist_ok=True)
STEM = "session-0-least-squares-foundations"
IMG_URL = "https://nurdauletakhanov.github.io/ReadingLab/slides/img/"

# ----------------------------------------------------------------------------- palette
BG = "#f7f7f5"; INK = "#24282e"; MUTED = "#6b7280"; ACCENT = "#a8710a"; ACCENT_SOFT = "#f6ecd6"
BLUE = "#37627f"; BLUE_SOFT = "#e8eff4"; LINE = "#d5d7d3"; RED = "#b0413a"

LECTURE = {"date": "Friday, 4 September 2026", "place": "Lecture Hall 1", "time": "6 pm"}
AUTHOR = "Nurdaulet Akhanov"

# ----------------------------------------------------------------------------- the data
cars = np.array([  # weight (t), horsepower, fuel (L/100 km)
    [1.2, 95, 6.4],
    [1.4, 100, 7.1],
    [1.0, 90, 5.8],
    [1.8, 150, 9.6],
    [1.6, 115, 8.2],
])
X = cars[:, :2]
y = cars[:, 2]
xbar = X.mean(axis=0)
ybar = y.mean()
Xc = X - xbar
yc = y - ybar
Xs = Xc / np.linalg.norm(Xc, axis=0)          # centred, unit length columns
R = Xs.T @ Xs                                  # the correlation matrix
r = R[0, 1]
lam, P = np.linalg.eigh(R)
order = np.argsort(lam)[::-1]
lam, P = lam[order], P[:, order]
P = P * np.where(P[0] < 0, -1, 1)          # draw eigenvectors pointing right
beta_c = np.linalg.solve(Xc.T @ Xc, Xc.T @ yc)  # slopes on centred data
beta0 = ybar - beta_c @ xbar
theta = np.degrees(np.arccos(r))
cov_wh = np.cov(X[:, 0], X[:, 1])[0, 1]           # tonne-horsepower
sd_w, sd_h = X.std(axis=0, ddof=1)


def f(x, d=2):
    return f"{x:.{d}f}"


# ----------------------------------------------------------------------------- LaTeX equations
TEX_RC = {"text.usetex": True, "text.latex.preamble": r"\usepackage{amsmath,amssymb}"}
EQ_DPI = 200


def render_tex(tex: str, name: str, size: int = 22) -> tuple[pathlib.Path, float, float]:
    """Typeset one display equation with LaTeX on the accent-soft panel. Returns (path, width_in, height_in)."""
    out = IMG / f"{name}.png"
    with plt.rc_context(TEX_RC):
        fig = plt.figure(figsize=(12, 3), dpi=EQ_DPI)
        fig.patch.set_facecolor(ACCENT_SOFT)
        fig.text(0.5, 0.5, f"${tex}$", ha="center", va="center", fontsize=size, color=INK)
        fig.savefig(out, facecolor=ACCENT_SOFT, bbox_inches="tight", pad_inches=0.22)
        plt.close(fig)
    from PIL import Image
    w, h = Image.open(out).size
    return out, w / EQ_DPI, h / EQ_DPI


# ----------------------------------------------------------------------------- figures
def style_axes(ax):
    ax.set_facecolor(BG)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    for s in ("left", "bottom"):
        ax.spines[s].set_color(LINE)
    ax.tick_params(colors=MUTED, labelsize=11)
    ax.xaxis.label.set_color(INK); ax.yaxis.label.set_color(INK)


def hull(points):
    pts = sorted(map(tuple, points))
    def half(seq):
        h = []
        for p in seq:
            while len(h) >= 2 and (h[-1][0]-h[-2][0])*(p[1]-h[-2][1]) - (h[-1][1]-h[-2][1])*(p[0]-h[-2][0]) <= 0:
                h.pop()
            h.append(p)
        return h
    lower, upper = half(pts), half(pts[::-1])
    return np.array(lower[:-1] + upper[:-1])


def fig_centering():
    fig, ax = plt.subplots(figsize=(9.2, 5.2), dpi=200)
    fig.patch.set_facecolor(BG); style_axes(ax)
    ax.axhline(0, color=LINE, lw=1); ax.axvline(0, color=LINE, lw=1)
    # the shift, drawn once per car so the eye sees five parallel arrows
    for (x0, y0), (x1, y1) in zip(X, Xc):
        ax.add_patch(FancyArrowPatch((x0, y0), (x1, y1), arrowstyle="-|>", mutation_scale=12,
                                     color=MUTED, lw=1, alpha=.55, shrinkA=6, shrinkB=6))
    for pts, col, soft in ((X, ACCENT, ACCENT_SOFT), (Xc, BLUE, BLUE_SOFT)):
        ax.add_patch(Polygon(hull(pts), closed=True, facecolor=soft, edgecolor=col, lw=1.4, alpha=.9, zorder=2))
        ax.scatter(pts[:, 0], pts[:, 1], s=70, color=col, zorder=4, edgecolor=BG, lw=1)
    ax.scatter(*xbar, marker="x", s=90, color=ACCENT, lw=2.2, zorder=5)
    ax.scatter(0, 0, marker="x", s=90, color=BLUE, lw=2.2, zorder=5)
    ax.annotate(f"original data\nmean ({xbar[0]:.1f}, {xbar[1]:.0f})", xy=xbar, xytext=(1.0, 150),
                ha="left", va="center", color=ACCENT, fontsize=12, fontweight="bold")
    ax.annotate("centred data\nmean (0, 0)", xy=(0, 0), xytext=(-0.95, -38),
                ha="left", va="center", color=BLUE, fontsize=12, fontweight="bold")
    ax.text(0.05, 105, f"every car moves by\n(−{xbar[0]:.1f}, −{xbar[1]:.0f})", ha="center", va="center",
            color=MUTED, fontsize=12, style="italic")
    ax.set_xlim(-1.05, 2.1); ax.set_ylim(-60, 165)
    ax.set_xlabel("weight (t)", fontsize=12); ax.set_ylabel("horsepower", fontsize=12)
    fig.tight_layout(); fig.savefig(IMG / "centering.png", facecolor=BG); plt.close(fig)


def fig_angle():
    fig, ax = plt.subplots(figsize=(5.4, 4.6), dpi=200)
    fig.patch.set_facecolor(BG); style_axes(ax)
    ax.set_aspect("equal"); ax.axis("off")
    t = np.radians(theta)
    v1 = np.array([1.0, 0.0]); v2 = np.array([np.cos(t), np.sin(t)])
    for v, col, lab in ((v1, ACCENT, "x₁  (weight)"), (v2, BLUE, "x₂  (horsepower)")):
        ax.add_patch(FancyArrowPatch((0, 0), v, arrowstyle="-|>", mutation_scale=18, color=col, lw=2.4))
        ax.text(*(v * 1.08), lab, color=col, fontsize=13, fontweight="bold",
                ha="left" if v is v1 else "center", va="center")
    arc = np.linspace(0, t, 40); ax.plot(0.35*np.cos(arc), 0.35*np.sin(arc), color=INK, lw=1.4)
    ax.text(0.42*np.cos(t/2), 0.42*np.sin(t/2), f"θ = {theta:.0f}°", fontsize=13, color=INK, ha="left", va="center")
    ax.text(0.55, -0.22, f"r = cos θ = {r:.3f}", fontsize=15, color=INK, ha="center")
    ax.set_xlim(-0.15, 1.75); ax.set_ylim(-0.35, 1.05)
    fig.tight_layout(); fig.savefig(IMG / "angle.png", facecolor=BG); plt.close(fig)


def fig_eigen():
    fig, ax = plt.subplots(figsize=(5.6, 5.2), dpi=200)
    fig.patch.set_facecolor(BG); style_axes(ax)
    ax.set_aspect("equal")
    ax.axhline(0, color=LINE, lw=1); ax.axvline(0, color=LINE, lw=1)
    ang = np.degrees(np.arctan2(P[1, 0], P[0, 0]))
    k = 1.15
    ax.add_patch(Ellipse((0, 0), 2*k*np.sqrt(lam[0]), 2*k*np.sqrt(lam[1]), angle=ang,
                         facecolor=ACCENT_SOFT, edgecolor=ACCENT, lw=1.2, alpha=.8))
    ax.scatter(Xs[:, 0], Xs[:, 1], s=70, color=INK, zorder=4, edgecolor=BG)
    for j, col in ((0, ACCENT), (1, BLUE)):
        v = P[:, j] * np.sqrt(lam[j]) * k
        ax.add_patch(FancyArrowPatch((0, 0), v, arrowstyle="-|>", mutation_scale=16, color=col, lw=2.4, zorder=5))
        off = np.array([0.06, 0.1]) if j == 0 else np.array([-0.9, 0.24])
        ax.text(*(v + off), f"λ{'₁' if j == 0 else '₂'} = {lam[j]:.2f}", color=col, fontsize=13, fontweight="bold")
    ax.set_xlim(-1.35, 1.35); ax.set_ylim(-1.25, 1.35)
    ax.set_xlabel("weight (centred, unit length)", fontsize=11); ax.set_ylabel("horsepower (centred, unit length)", fontsize=11)
    fig.tight_layout(); fig.savefig(IMG / "eigen.png", facecolor=BG); plt.close(fig)


def fig_projection():
    fig, ax = plt.subplots(figsize=(6.4, 4.6), dpi=200)
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG); ax.axis("off")
    plane = np.array([[0.15, 1.05], [5.6, 0.25], [7.5, 1.95], [2.05, 2.75]])
    ax.add_patch(Polygon(plane, closed=True, facecolor=BLUE_SOFT, edgecolor=BLUE, lw=1.4, alpha=.95))
    O = np.array([1.75, 1.45])
    yhat = np.array([5.15, 1.62])
    Y = np.array([5.15, 4.05])
    for tip, lab, col in (((4.35, 1.06), "x₁", MUTED), ((3.35, 2.25), "x₂", MUTED)):
        ax.add_patch(FancyArrowPatch(O, tip, arrowstyle="-|>", mutation_scale=13, color=col, lw=1.6))
        ax.text(*(np.array(tip) + (0.12, -0.16)), lab, color=col, fontsize=13, fontweight="bold")
    ax.add_patch(FancyArrowPatch(O, Y, arrowstyle="-|>", mutation_scale=17, color=ACCENT, lw=2.6, zorder=5))
    ax.add_patch(FancyArrowPatch(O, yhat, arrowstyle="-|>", mutation_scale=17, color=BLUE, lw=2.6, zorder=5))
    ax.add_patch(FancyArrowPatch(yhat, Y, arrowstyle="-|>", mutation_scale=15, color=RED, lw=2.2,
                                 linestyle=(0, (5, 3)), zorder=5))
    d = 0.3
    ax.plot([yhat[0], yhat[0] + d, yhat[0] + d], [yhat[1] + d, yhat[1] + d, yhat[1]],
            color=RED, lw=1.3, zorder=6)
    ax.scatter(*O, s=28, color=INK, zorder=7)
    ax.text(O[0] - 0.3, O[1] - 0.1, "0", color=INK, fontsize=12)
    ax.text(Y[0] + 0.15, Y[1], "Y", color=ACCENT, fontsize=16, fontweight="bold", va="center")
    ax.text(yhat[0] + 0.2, yhat[1] - 0.28, "Ŷ = Xβ̂", color=BLUE, fontsize=15, fontweight="bold", va="top")
    ax.text(Y[0] + 0.2, (Y[1] + yhat[1]) / 2, "e = Y − Ŷ", color=RED, fontsize=14, fontweight="bold", va="center")
    ax.text(6.75, 2.25, "Col(X)", color=BLUE, fontsize=14, fontweight="bold", ha="center")
    ax.set_xlim(-0.35, 8.6); ax.set_ylim(0.0, 4.6)
    fig.tight_layout(); fig.savefig(IMG / "projection.png", facecolor=BG); plt.close(fig)


def fig_dartboard():
    rng = np.random.default_rng(7)
    fig, axes = plt.subplots(1, 2, figsize=(9.2, 4.6), dpi=200)
    fig.patch.set_facecolor(BG)
    specs = (("Unbiased, high variance", (0, 0), 0.55, ACCENT), ("Biased, low variance", (0.28, 0.2), 0.14, BLUE))
    for ax, (title, mu, sd, col) in zip(axes, specs):
        ax.set_facecolor(BG); ax.set_aspect("equal"); ax.axis("off")
        for rad in (1.0, 0.7, 0.4, 0.15):
            ax.add_patch(plt.Circle((0, 0), rad, facecolor="none", edgecolor=LINE, lw=1.3))
        ax.scatter(0, 0, marker="+", s=140, color=INK, lw=1.5)
        pts = rng.normal(mu, sd, size=(14, 2))
        ax.scatter(pts[:, 0], pts[:, 1], s=55, color=col, edgecolor=BG, zorder=3)
        ax.scatter(*mu, marker="x", s=90, color=col, lw=2.2, zorder=4)
        ax.set_title(title, color=INK, fontsize=14, pad=10)
        ax.set_xlim(-1.35, 1.35); ax.set_ylim(-1.3, 1.3)
    fig.tight_layout(rect=(0, 0, 1, 0.97)); fig.savefig(IMG / "dartboard.png", facecolor=BG); plt.close(fig)


# ----------------------------------------------------------------------------- content
# Slide kinds: title, table, text (lines + optional equation), figure (image + lines), columns, checklist.
# Lines are short sentences the lecturer can say; equations use Unicode so PowerPoint needs no add-ins.
SLIDES = [
    dict(kind="title", title="Least Squares Foundations",
         sub="Session 0 · what the ridge paper assumes you already know",
         meta=[AUTHOR, f"{LECTURE['date']} · {LECTURE['time']} · {LECTURE['place']}"]),

    dict(kind="text", label="Why a Session 0", title="Fourteen pages that take five things for granted",
         lines=["Hoerl & Kennard (1970) never explain these. Today we build them from one table of cars:"],
         bullets=["the linear model written as  Y = Xβ + ε",
                  "correlation between predictors as an angle",
                  "the eigenvalues of XᵀX, and why a small one is bad news",
                  "least squares: unbiased, with covariance σ²(XᵀX)⁻¹",
                  "error = variance + bias², so unbiased is not the same as accurate"],
         note="Session 1 uses all five to break least squares, then to fix it."),

    dict(kind="table", label="Part I · The model", title="Five cars, three numbers each",
         header=["Car", "Weight (t)", "Horsepower", "Fuel (L/100 km)"],
         rows=[[str(i + 1), f"{w:.1f}", f"{h:.0f}", f"{fu:.1f}"] for i, (w, h, fu) in enumerate(cars)],
         lines=["We want to predict fuel from the other two columns.",
                "Fuel is the response. Weight and horsepower are the predictors."]),

    dict(kind="text", label="Part I · The model", title="From the table to an equation",
         equation=r"\text{fuel}_i = \beta_1\,\text{weight}_i + \beta_2\,\text{horsepower}_i + \varepsilon_i",
         bullets=["β₁ and β₂ belong to the relationship, not to any single car. They are what we want.",
                  "εᵢ is everything the two measurements miss: driver, tyres, weather, the scale.",
                  "No amount of staring at the table reveals β. We will have to estimate it."]),

    dict(kind="text", label="Part I · The model", title="Five equations become one",
         equation=r"\underbrace{\begin{pmatrix}6.4\\7.1\\5.8\\9.6\\8.2\end{pmatrix}}_{Y} = \underbrace{\begin{pmatrix}1.2&95\\1.4&100\\1.0&90\\1.8&150\\1.6&115\end{pmatrix}}_{X} \underbrace{\begin{pmatrix}\beta_1\\ \beta_2\end{pmatrix}}_{\beta} + \underbrace{\begin{pmatrix}\varepsilon_1\\ \varepsilon_2\\ \varepsilon_3\\ \varepsilon_4\\ \varepsilon_5\end{pmatrix}}_{\varepsilon}",
         bullets=["Y = Xβ + ε.   X has n = 5 rows, one per car, and p = 2 columns, one per predictor.",
                  "Every paper on the reading list starts from this line."]),

    dict(kind="text", label="Part I · The model", title="The intercept is a question about the origin",
         equation=r"y_i = \beta_0 + \beta_1 x_{i1} + \beta_2 x_{i2} + \varepsilon_i",
         bullets=["A real model has an intercept β₀; it lets the fitted plane sit above or below zero.",
                  "Without β₀ the fit must pass through the origin: a car with zero weight and zero horsepower.",
                  "Carrying β₀ through every formula is a nuisance. The cleaner move is to relocate the origin."]),

    dict(kind="figure", label="Part I · The model", title="Centring moves the origin to the average car",
         image="centering.png", width=7.3,
         equation=rf"\bar x = ({xbar[0]:.1f},\,{xbar[1]:.0f}), \qquad x_{{ij}} \mapsto x_{{ij}} - \bar x_j",
         bullets=["Subtract each column's mean from every entry, and do the same to y.",
                  "Every car moves by the same vector, so the cloud slides as one rigid piece.",
                  "Distances, angles, and the shape are untouched. Only the mean has moved, to (0, 0)."]),

    dict(kind="text", label="Part I · The model", title="What centring buys us",
         equation=r"\hat\beta_0 = \bar y - \hat\beta_1 \bar x_1 - \hat\beta_2 \bar x_2",
         bullets=["The average car is now the origin, and its predicted fuel is the average fuel: the intercept is exactly zero.",
                  "So we fit only the slopes, on centred data, from  Y = Xβ + ε  with no β₀ in sight.",
                  f"The intercept comes back for free afterwards. For the cars: β̂₁ = {f(beta_c[0])}, β̂₂ = {f(beta_c[1],3)}, β̂₀ = {f(beta0)}.",
                  "Ridge regression shrinks slopes, never the intercept. Centring is what keeps the two apart."]),

    dict(kind="text", label="Part II · The geometry of correlation", title="Predictors are vectors",
         equation=r"X = [\,x_1 \;\; x_2\,], \quad x_1, x_2 \in \mathbb{R}^5, \qquad X^{\top}X = \begin{pmatrix} x_1^{\top}x_1 & x_1^{\top}x_2 \\ x_2^{\top}x_1 & x_2^{\top}x_2 \end{pmatrix}",
         bullets=["Read X by columns: weight is one vector in ℝ⁵, horsepower is another.",
                  "XᵀX collects their dot products. It is p × p, so rows and columns now index predictors, not cars.",
                  "The whole story of this session and the next lives inside that small matrix."]),

    dict(kind="figure", label="Part II · The geometry of correlation", title="Correlation is an angle",
         image="angle.png", width=4.9, side=True,
         equation=r"x_1^{\top}x_2 = \|x_1\|\,\|x_2\| \cos\theta = \cos\theta = r",
         bullets=["Centre both columns and scale each to unit length.",
                  f"Their dot product is the cosine of the angle between them, and that is the sample correlation. For the cars, r = {r:.3f}.",
                  "|r| ≤ 1 is Cauchy–Schwarz, nothing more.",
                  "Nearly parallel columns mean nearly redundant predictors."]),

    dict(kind="text", label="Part II · The geometry of correlation", title="Why standardize? A detour through covariance",
         equation=r"\operatorname{cov}(u,v) = \frac{1}{n-1}\sum_{i=1}^{n}(u_i-\bar u)(v_i-\bar v), \qquad r = \frac{\operatorname{cov}(u,v)}{\operatorname{sd}(u)\,\operatorname{sd}(v)}",
         size=17,
         bullets=[f"Covariance is the crude measure of co-movement: centre, multiply entry by entry, average. For the cars, cov(weight, hp) = {cov_wh:.1f}.",
                  "But 7.0 of what? Tonne-horsepower. Measure weight in grams and the same five cars give 7 000 000. The number carries the ruler.",
                  f"The standard deviation is the typical distance from the mean, in the variable's own units: sd(u) = √cov(u, u). Here sd(weight) = {sd_w:.4f} t and sd(hp) = {sd_h:.3f}.",
                  f"Divide by the standard deviations and the ruler cancels: {cov_wh:.1f} / ({sd_w:.4f} × {sd_h:.3f}) = {r:.3f}, in tonnes or in grams alike.",
                  "Standardizing does that division once, on the data itself. Afterwards a plain dot product already is the correlation, and the ridge paper assumes exactly that."]),

    dict(kind="text", label="Part II · The geometry of correlation", title="XᵀX is the correlation matrix",
         equation=rf"X^{{\top}}X = \begin{{pmatrix}} 1 & r \\ r & 1 \end{{pmatrix}} = \begin{{pmatrix}} 1 & {r:.3f} \\ {r:.3f} & 1 \end{{pmatrix}}",
         bullets=["On the diagonal: squared column lengths, all equal to 1.",
                  "Off the diagonal: the pairwise correlations.",
                  "The trace is p, so the eigenvalues always add up to p. They share a fixed budget."]),

    dict(kind="figure", label="Part II · The geometry of correlation", title="Eigenvalues: where the information is",
         image="eigen.png", width=4.9, side=True,
         equation=r"X^{\top}X = P\Lambda P^{\top}, \qquad X^{\top}X\,p_j = \lambda_j\,p_j",
         bullets=["Each eigenvalue is the spread of the data cloud along its eigenvector.",
                  f"For the cars: λ₁ = {lam[0]:.2f} along the diagonal, λ₂ = {lam[1]:.2f} across it. They sum to 2.",
                  "Along the short axis the data barely varies. Whatever the coefficients do in that direction, the data cannot tell.",
                  "Remember the small one. Everything in Session 1 is about 1/λ₂."]),

    dict(kind="text", label="Part II · The geometry of correlation", title="What we assume about the noise",
         equation=r"\mathbb{E}[\varepsilon \mid X] = 0, \qquad \operatorname{Cov}(\varepsilon \mid X) = \sigma^2 I",
         bullets=["Nothing systematic is left over: on average the model is right.",
                  "Every car has the same error variance, and no two errors are linked.",
                  "That is all. Normality is not needed for anything today."]),

    dict(kind="text", label="Part III · What makes an estimate good", title="Least squares: make the leftovers small",
         equation=r"\hat\beta = \arg\min_{\beta} \|Y - X\beta\|^2 \;\Longrightarrow\; X^{\top}X\hat\beta = X^{\top}Y \;\Longrightarrow\; \hat\beta = (X^{\top}X)^{-1}X^{\top}Y",
         bullets=["The objective is the squared length of the residual vector in ℝ⁵.",
                  "Its gradient is −2XᵀY + 2XᵀXβ. Setting it to zero gives the normal equations.",
                  "If XᵀX can be inverted, the solution is unique and explicit.",
                  "Everything about β̂ is decided by that inverse."]),

    dict(kind="figure", label="Part III · What makes an estimate good", title="Least squares as orthogonal projection",
         image="projection.png", width=6.1, side=True,
         equation=r"\hat Y = X\hat\beta \in \operatorname{Col}(X), \qquad e = Y - \hat Y \perp \operatorname{Col}(X)",
         bullets=["Every fitted vector has the form XB, so it lies in the column space of X.",
                  "Least squares selects Ŷ = Xβ̂, the point in Col(X) nearest to Y.",
                  "At the nearest point, the residual e = Y − Ŷ is perpendicular to every column of X. Thus Xᵀe = 0, which gives the normal equations.",
                  "Ŷ lives in ℝⁿ. β̂ lives in ℝᵖ and gives the coordinates of Ŷ in the columns of X."],
         note="This derivation uses geometry. Differentiating the squared residual length gives the same condition."),

    dict(kind="text", label="Part III · What makes an estimate good", title="When can XᵀX be inverted?",
         equation=r"v^{\top}X^{\top}Xv = (Xv)^{\top}(Xv) = \|Xv\|^2 \ge 0",
         bullets=["So XᵀX is positive semidefinite for any X.",
                  "If the columns are linearly independent, Xv ≠ 0 whenever v ≠ 0, the quadratic form is strictly positive, and the inverse exists.",
                  "Nearly dependent columns keep the inverse but make it enormous. That is the ridge paper's whole subject."]),

    dict(kind="text", label="Part III · What makes an estimate good", title="The estimate is a random variable",
         equation=r"\hat\beta = (X^{\top}X)^{-1}X^{\top}(X\beta + \varepsilon) = \beta + (X^{\top}X)^{-1}X^{\top}\varepsilon",
         bullets=["β and X are fixed. Redraw the noise, and the estimate moves.",
                  "An estimator therefore has a distribution, and we can ask two questions of it:",
                  "where is its centre, and how wide is its spread?"]),

    dict(kind="text", label="Part III · What makes an estimate good", title="Unbiased, with covariance σ²(XᵀX)⁻¹",
         equation=r"\mathbb{E}[\hat\beta \mid X] = \beta, \qquad \operatorname{Cov}(\hat\beta \mid X) = \sigma^2 (X^{\top}X)^{-1}",
         bullets=["The centre is exactly right: least squares is unbiased.",
                  "The spread is the inverse of the correlation matrix. Along eigenvector pⱼ the variance is σ²/λⱼ.",
                  f"For the cars: σ²/λ₁ = {1/lam[0]:.2f} σ²  but  σ²/λ₂ = {1/lam[1]:.1f} σ².",
                  f"One direction of the estimate is {lam[0]/lam[1]:.0f} times noisier than the other, and the data caused it."]),

    dict(kind="figure", label="Part III · What makes an estimate good", title="Unbiased is not the same as accurate",
         image="dartboard.png", width=7.3,
         equation=r"\operatorname{MSE}(\tilde\beta) = \mathbb{E}\|\tilde\beta - \beta\|^2 = \operatorname{tr}\operatorname{Cov}(\tilde\beta) + \|\operatorname{bias}(\tilde\beta)\|^2",
         bullets=["Player A is centred on the bullseye but scatters widely. Player B aims a little off and clusters tightly.",
                  "Every single dart of B lands closer. Mean squared error scores both players on the same scale."]),

    dict(kind="text", label="Part III · What makes an estimate good", title="Gauss–Markov, read carefully",
         equation=r"\operatorname{Cov}(\tilde\beta \mid X) - \operatorname{Cov}(\hat\beta \mid X) \succeq 0 \quad \text{for every } \tilde\beta \text{ linear in } Y \text{ and unbiased}",
         bullets=["Among linear unbiased estimators, least squares has the smallest covariance.",
                  "Two qualifiers. Drop “unbiased” and the theorem says nothing at all.",
                  "That gap is the door ridge regression walks through."]),

    dict(kind="text", label="Part III · What makes an estimate good", title="Where this leaves us",
         equation=r"\hat\beta = (X^{\top}X)^{-1}X^{\top}Y \qquad\qquad \hat\beta(k) = (X^{\top}X + kI)^{-1}X^{\top}Y, \quad k > 0",
         bullets=["In the eigenvector basis, every 1/λⱼ becomes 1/(λⱼ + k). The small eigenvalue can no longer explode.",
                  "The price is bias. Session 1 proves that for some k > 0 the trade lowers the mean squared error, always.",
                  "Bring the centred model, the eigenvalues, and the bias–variance split. The paper assumes all three."]),

    dict(kind="columns", label="Closing", title="Exercises",
         left=["1.  Show that XᵀX is positive semidefinite.",
               "2.  Show it is invertible exactly when the columns of X are linearly independent.",
               "3.  For centred, unit-length columns, show that the dot product is the correlation and lies in [−1, 1].",
               "4.  Show that the eigenvalues of XᵀX sum to p under that scaling."],
         right=["5.  Prove the projection condition, then derive the normal equations and the least squares solution.",
                "6.  Prove that least squares is unbiased and derive its covariance.",
                "7.  Prove  E‖β̃ − β‖² = tr Cov(β̃) + ‖bias‖².",
                "8.  For β̃ = cβ̂ with 0 ≤ c ≤ 1, find the c that minimises the MSE. Is it 1?"],
         note="Work in pairs. Mark the line where each assumption about ε is used."),

    dict(kind="checklist", label="Closing", title="Before Session 1 you should be able to prove",
         items=["XᵀX is positive semidefinite, and positive definite under full column rank.",
                "For centred unit-length columns, a dot product is a correlation.",
                "The eigenvalues of XᵀX sum to p.",
                "Orthogonal projection gives the normal equations. Full column rank makes the least squares coefficients unique.",
                "Least squares is unbiased with covariance σ²(XᵀX)⁻¹.",
                "Vector MSE = total variance + squared bias.",
                "Gauss–Markov compares only linear unbiased estimators.",
                "A biased estimator can have a lower MSE."],
         note="Next: Hoerl & Kennard (1970), “Ridge Regression: Biased Estimation for Nonorthogonal Problems”."),
]

# ----------------------------------------------------------------------------- pptx
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.75)


def rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#"))


def add_text(slide, x, y, w, h, text, size=20, bold=False, color=INK, font="Arial", align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = text if isinstance(text, list) else [text]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = para
        run.font.size = Pt(size); run.font.bold = bold; run.font.name = font; run.font.color.rgb = rgb(color)
    return box


def bullets_box(slide, x, y, w, h, items, size=19):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.18; p.space_after = Pt(9)
        dot = p.add_run(); dot.text = "•  "; dot.font.size = Pt(size); dot.font.color.rgb = rgb(ACCENT); dot.font.name = "Arial"; dot.font.bold = True
        run = p.add_run(); run.text = item; run.font.size = Pt(size); run.font.color.rgb = rgb(INK); run.font.name = "Arial"
    return box


def chrome(slide, label, title, number):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(BG)
    add_text(slide, MARGIN, Inches(0.42), W - 2*MARGIN, Inches(0.3), label.upper(), size=12, bold=True, color=ACCENT)
    add_text(slide, MARGIN, Inches(0.72), W - 2*MARGIN, Inches(0.9), title, size=34, bold=True, color=INK, font="Georgia")
    add_text(slide, W - MARGIN - Inches(1), H - Inches(0.55), Inches(1), Inches(0.3), str(number), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def equation_box(slide, y, tex, mono=False, size=None):  # noqa: E302
    return equation_box_width(slide, MARGIN, y, W - 2*MARGIN, tex)




def build_pptx():
    prs = Presentation(); prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]
    for n, s in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        k = s["kind"]
        if k == "title":
            bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(BG)
            add_text(slide, MARGIN, Inches(1.2), W - 2*MARGIN, Inches(0.3), "READINGLAB · SESSION 0", size=13, bold=True, color=ACCENT)
            add_text(slide, MARGIN, Inches(1.65), Inches(10), Inches(2.2), s["title"], size=60, bold=True, color=INK, font="Georgia")
            add_text(slide, MARGIN, Inches(3.8), Inches(10.5), Inches(0.6), s["sub"], size=24, color=MUTED, font="Georgia")
            add_text(slide, MARGIN, Inches(5.4), Inches(10), Inches(1.0), s["meta"], size=18, color=INK, line_spacing=1.4)
            continue
        chrome(slide, s["label"], s["title"], n)
        y = Inches(1.75)
        if k == "text":
            if s.get("lines"):
                add_text(slide, MARGIN, y, W - 2*MARGIN, Inches(0.9), s["lines"], size=20, line_spacing=1.25, space_after=6)
                y += Inches(0.55) * len(s["lines"]) + Inches(0.2)
            if s.get("equation"):
                y = equation_box(slide, y, s["equation"]) + Inches(0.35)
            if s.get("bullets"):
                bullets_box(slide, MARGIN, y, W - 2*MARGIN, H - y - Inches(0.9), s["bullets"], size=s.get("size", 19))
            if s.get("note"):
                add_text(slide, MARGIN, H - Inches(1.0), W - 2*MARGIN, Inches(0.4), s["note"], size=16, color=MUTED, font="Georgia")
        elif k == "table":
            rows, cols = len(s["rows"]) + 1, len(s["header"])
            tw = Inches(7.2); th = Inches(0.46) * rows
            shape = slide.shapes.add_table(rows, cols, MARGIN, y, tw, th)
            tbl = shape.table
            for c, name in enumerate(s["header"]):
                cell = tbl.cell(0, c); cell.text = name
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(ACCENT_SOFT)
            for r_, row in enumerate(s["rows"], start=1):
                for c, val in enumerate(row):
                    cell = tbl.cell(r_, c); cell.text = val
                    cell.fill.solid(); cell.fill.fore_color.rgb = rgb("#ffffff" if r_ % 2 else "#fbfaf6")
            for r_ in range(rows):
                for c in range(cols):
                    for p in tbl.cell(r_, c).text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                        for run in p.runs:
                            run.font.size = Pt(18); run.font.name = "Arial"; run.font.color.rgb = rgb(INK); run.font.bold = r_ == 0
            add_text(slide, MARGIN + tw + Inches(0.6), y + Inches(0.2), W - 2*MARGIN - tw - Inches(0.6), Inches(3),
                     s["lines"], size=20, line_spacing=1.3, space_after=12)
        elif k == "figure":
            img = IMG / s["image"]
            if s.get("top"):
                pic_w = Inches(s["width"])
                pic = slide.shapes.add_picture(str(img), int((W - pic_w) / 2), y, width=pic_w)
                y2 = y + pic.height + Inches(0.35)
                y2 = equation_box(slide, y2, s["equation"]) + Inches(0.3)
                bullets_box(slide, MARGIN, y2, W - 2*MARGIN, H - y2 - Inches(0.6), s["bullets"], size=18)
            elif s.get("side"):
                pic_w = Inches(s["width"])
                slide.shapes.add_picture(str(img), W - MARGIN - pic_w, y, width=pic_w)
                tw = W - 2*MARGIN - pic_w - Inches(0.5)
                y = equation_box_width(slide, MARGIN, y, tw, s["equation"], size=20) + Inches(0.3)
                bullets_box(slide, MARGIN, y, tw, H - y - Inches(0.8), s["bullets"], size=18)
            else:
                pic_w = Inches(s["width"])
                pic = slide.shapes.add_picture(str(img), MARGIN, y, width=pic_w)
                tw = W - 2*MARGIN - pic_w - Inches(0.45)
                x = MARGIN + pic_w + Inches(0.45)
                y2 = equation_box_width(slide, x, y, tw, s["equation"], size=18) + Inches(0.25)
                bullets_box(slide, x, y2, tw, H - y2 - Inches(0.7), s["bullets"], size=16)
        elif k == "columns":
            cw = (W - 2*MARGIN - Inches(0.6)) / 2
            add_text(slide, MARGIN, y, cw, Inches(4.5), s["left"], size=18, line_spacing=1.25, space_after=12)
            add_text(slide, MARGIN + cw + Inches(0.6), y, cw, Inches(4.5), s["right"], size=18, line_spacing=1.25, space_after=12)
            add_text(slide, MARGIN, H - Inches(1.0), W - 2*MARGIN, Inches(0.4), s["note"], size=16, color=MUTED, font="Georgia")
        elif k == "checklist":
            items = [f"{i}.  {t}" for i, t in enumerate(s["items"], start=1)]
            half = (len(items) + 1) // 2
            cw = (W - 2*MARGIN - Inches(0.6)) / 2
            add_text(slide, MARGIN, y, cw, Inches(4.5), items[:half], size=18, line_spacing=1.25, space_after=12)
            add_text(slide, MARGIN + cw + Inches(0.6), y, cw, Inches(4.5), items[half:], size=18, line_spacing=1.25, space_after=12)
            add_text(slide, MARGIN, H - Inches(1.0), W - 2*MARGIN, Inches(0.4), s["note"], size=16, color=MUTED, font="Georgia")
    prs.save(OUT / f"{STEM}.pptx")


_eq_counter = [0]


def equation_box_width(slide, x, y, w, tex, size=22):
    _eq_counter[0] += 1
    path, win, hin = render_tex(tex, f"eq{_eq_counter[0]:02d}", size)
    max_h = 2.6
    scale = min(1.0, (w / 914400) / win, max_h / hin)
    pw, ph = Inches(win * scale), Inches(hin * scale)
    slide.shapes.add_picture(str(path), int(x + (w - pw) / 2), y, width=pw, height=ph)
    return y + ph


# ----------------------------------------------------------------------------- marp markdown
MARP_HEAD = f"""---
marp: true
theme: default
paginate: true
math: mathjax
size: 16:9
style: |
  section {{ background: {BG}; color: {INK}; font-family: Arial, Helvetica, sans-serif; font-size: 26px; line-height: 1.35; padding: 54px 72px; }}
  h1, h2 {{ color: {INK}; font-family: Georgia, "Times New Roman", serif; font-weight: 600; }}
  h1 {{ font-size: 60px; line-height: 1.08; }}
  h2 {{ font-size: 40px; margin: 0 0 22px; }}
  table {{ font-size: 22px; }}
  th {{ background: {ACCENT_SOFT}; }}
  ul {{ padding-left: 1.1em; }}
  li {{ margin-bottom: .35em; }}
  li::marker {{ color: {ACCENT}; }}
  .label {{ color: {ACCENT}; font-size: 15px; font-weight: 700; letter-spacing: .06em; text-transform: uppercase; margin: 0 0 8px; }}
  .eq {{ background: {ACCENT_SOFT}; border-radius: 8px; padding: 6px 20px; text-align: center; font-size: 26px; margin: 12px 0 22px; }}
  .eq p {{ margin: 0; }}
  .mono {{ font-family: Menlo, monospace; font-size: 20px; text-align: left; display: inline-block; }}
  .note {{ color: {MUTED}; font-family: Georgia, serif; font-size: 20px; margin-top: 26px; }}
  .subtitle {{ color: {MUTED}; font-family: Georgia, serif; font-size: 28px; }}
  .side {{ display: grid; grid-template-columns: 1fr 420px; gap: 36px; align-items: start; }}
  .wide {{ display: grid; grid-template-columns: 640px 1fr; gap: 32px; align-items: start; }}
  .two {{ display: grid; grid-template-columns: 1fr 1fr; gap: 48px; }}
  img {{ background: transparent; }}
---
"""


def md_eq(tex, mono=False):
    return '<div class="eq">\n\n$$\n' + tex + '\n$$\n\n</div>\n\n'


def build_md():
    out = [MARP_HEAD]
    for s in SLIDES:
        k = s["kind"]
        if k == "title":
            out.append(f'\n<p class="label">ReadingLab · Session 0</p>\n\n# {s["title"]}\n\n<p class="subtitle">{s["sub"]}</p>\n\n'
                       + "<br>".join(s["meta"]) + "\n")
            continue
        out.append(f'\n---\n\n<p class="label">{s["label"]}</p>\n\n## {s["title"]}\n\n')
        if k == "text":
            for l in s.get("lines", []):
                out.append(l + "\n\n")
            if s.get("equation"):
                out.append(md_eq(s["equation"], s.get("mono", False)))
            for b in s.get("bullets", []):
                out.append(f"- {b}\n")
            if s.get("note"):
                out.append(f'\n<p class="note">{s["note"]}</p>\n')
        elif k == "table":
            out.append("| " + " | ".join(s["header"]) + " |\n|" + "|".join("---:" if i else "---" for i in range(len(s["header"]))) + "|\n")
            for row in s["rows"]:
                out.append("| " + " | ".join(row) + " |\n")
            out.append("\n" + "\n\n".join(s["lines"]) + "\n")
        elif k == "figure":
            cls = "side" if s.get("side") else "wide"
            img = f'![]({IMG_URL + s["image"]})'
            body = md_eq(s["equation"]) + "".join(f"- {b}\n" for b in s["bullets"])
            if s.get("top"):
                out.append(f'<p style="text-align:center">{img.replace("![]", "![w:720]")}</p>\n\n{body}')
            elif s.get("side"):
                out.append(f'<div class="{cls}">\n<div>\n\n{body}\n</div>\n<div>\n\n{img}\n\n</div>\n</div>\n')
            else:
                out.append(f'<div class="{cls}">\n<div>\n\n{img}\n\n</div>\n<div>\n\n{body}\n</div>\n</div>\n')
        elif k == "columns":
            out.append('<div class="two">\n<div>\n\n' + "\n\n".join(s["left"]) + '\n\n</div>\n<div>\n\n' + "\n\n".join(s["right"]) + "\n\n</div>\n</div>\n")
            out.append(f'\n<p class="note">{s["note"]}</p>\n')
        elif k == "checklist":
            items = [f"{i}. {t}" for i, t in enumerate(s["items"], start=1)]
            half = (len(items) + 1) // 2
            out.append('<div class="two">\n<div>\n\n' + "\n".join(items[:half]) + '\n\n</div>\n<div>\n\n' + "\n".join(items[half:]) + "\n\n</div>\n</div>\n")
            out.append(f'\n<p class="note">{s["note"]}</p>\n')
    (OUT / f"{STEM}.md").write_text("".join(out))


if __name__ == "__main__":
    fig_centering(); fig_angle(); fig_eigen(); fig_dartboard(); fig_projection()
    build_pptx(); build_md()
    print(f"r = {r:.3f}, theta = {theta:.1f} deg, lambda = {lam.round(3)}, beta_c = {beta_c.round(3)}, beta0 = {beta0:.3f}")
    print(f"wrote {OUT / (STEM + '.pptx')} ({(OUT / (STEM + '.pptx')).stat().st_size/1e3:.0f} KB), {len(SLIDES)} slides")
